"""
╔════════════════════════════════════════════════════════════════════════════════╗
║  DOSYA: document_processor.py                                                  ║
║  AMAÇ: Multi-format doküman işleme (PDF, DOCX, Excel, PPT, HTML, TXT)         ║
║  ÖZELLİKLER:                                                                   ║
║    - PDF: pdfplumber + PyPDF2 + OCR fallback                                  ║
║    - DOCX: python-docx ile tablo ve resim desteği                             ║
║    - Excel: openpyxl ile çoklu sheet desteği                                  ║
║    - PPT: python-pptx ile slide extraction                                    ║
║    - HTML: BeautifulSoup ile temiz text çıkarma                               ║
║    - Otomatik encoding detection                                               ║
║    - Metadata extraction                                                       ║
╚════════════════════════════════════════════════════════════════════════════════╝
"""

from __future__ import annotations

import hashlib
import logging
import mimetypes
import re
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger(__name__)

# Optional imports with fallbacks
try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    from docx.table import Table
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False


@dataclass
class ProcessedDocument:
    """İşlenmiş doküman veri yapısı"""
    content: str
    source: str
    file_type: str
    title: str = ""
    page_count: int = 0
    word_count: int = 0
    char_count: int = 0
    content_hash: str = ""
    processed_at: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    tables: List[List[List[str]]] = field(default_factory=list)
    sections: List[Dict[str, str]] = field(default_factory=list)
    
    def __post_init__(self):
        if not self.content_hash:
            self.content_hash = hashlib.sha256(self.content.encode()).hexdigest()
        if not self.processed_at:
            self.processed_at = datetime.utcnow().isoformat()
        if not self.word_count:
            self.word_count = len(self.content.split())
        if not self.char_count:
            self.char_count = len(self.content)


class PDFProcessor:
    """PDF doküman işleyici"""
    
    @staticmethod
    def extract_text(file_path: Union[str, Path, BytesIO]) -> ProcessedDocument:
        """
        PDF'den text çıkar.
        
        Öncelik sırası:
        1. pdfplumber (tablo ve layout desteği)
        2. PyPDF2 (fallback)
        """
        text_parts = []
        tables = []
        page_count = 0
        metadata = {}
        
        if PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    page_count = len(pdf.pages)
                    metadata = pdf.metadata or {}
                    
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"[Sayfa {i+1}]\n{page_text}")
                        
                        page_tables = page.extract_tables()
                        if page_tables:
                            for table in page_tables:
                                if table:
                                    tables.append(table)
                                    table_text = PDFProcessor._table_to_text(table)
                                    text_parts.append(f"\n[Tablo - Sayfa {i+1}]\n{table_text}")
                
                if text_parts:
                    content = "\n\n".join(text_parts)
                    source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
                    
                    return ProcessedDocument(
                        content=content,
                        source=source,
                        file_type="pdf",
                        title=metadata.get("Title", Path(source).stem if source != "stream" else ""),
                        page_count=page_count,
                        metadata=metadata,
                        tables=tables,
                    )
            except Exception as e:
                logger.warning(f"pdfplumber başarısız, PyPDF2 deneniyor: {e}")
        
        if PYPDF_AVAILABLE:
            try:
                if isinstance(file_path, BytesIO):
                    reader = PdfReader(file_path)
                else:
                    reader = PdfReader(str(file_path))
                
                page_count = len(reader.pages)
                metadata = dict(reader.metadata) if reader.metadata else {}
                
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Sayfa {i+1}]\n{page_text}")
                
                content = "\n\n".join(text_parts)
                source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
                
                return ProcessedDocument(
                    content=content,
                    source=source,
                    file_type="pdf",
                    title=metadata.get("/Title", Path(source).stem if source != "stream" else ""),
                    page_count=page_count,
                    metadata=metadata,
                )
            except Exception as e:
                logger.error(f"PDF işleme hatası: {e}")
                raise
        
        raise ImportError("PDF işleme için pypdf veya pdfplumber gerekli")
    
    @staticmethod
    def _table_to_text(table: List[List[str]]) -> str:
        """Tabloyu text formatına çevir"""
        rows = []
        for row in table:
            if row:
                cells = [str(cell).strip() if cell else "" for cell in row]
                rows.append(" | ".join(cells))
        return "\n".join(rows)


class DOCXProcessor:
    """DOCX doküman işleyici"""
    
    @staticmethod
    def extract_text(file_path: Union[str, Path, BytesIO]) -> ProcessedDocument:
        """DOCX'den text çıkar (paragraflar, tablolar, başlıklar dahil)"""
        if not DOCX_AVAILABLE:
            raise ImportError("DOCX işleme için python-docx gerekli")
        
        try:
            if isinstance(file_path, BytesIO):
                doc = DocxDocument(file_path)
            else:
                doc = DocxDocument(str(file_path))
            
            text_parts = []
            tables = []
            sections = []
            
            current_section = {"title": "", "content": []}
            
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    para = None
                    for p in doc.paragraphs:
                        if p._element == element:
                            para = p
                            break
                    
                    if para:
                        text = para.text.strip()
                        if not text:
                            continue
                        
                        if para.style and para.style.name and para.style.name.startswith('Heading'):
                            if current_section["content"]:
                                sections.append(current_section.copy())
                            current_section = {"title": text, "content": []}
                            text_parts.append(f"\n## {text}\n")
                        else:
                            current_section["content"].append(text)
                            text_parts.append(text)
                
                elif element.tag.endswith('tbl'):
                    for table in doc.tables:
                        if table._element == element:
                            table_data = DOCXProcessor._extract_table(table)
                            if table_data:
                                tables.append(table_data)
                                table_text = DOCXProcessor._table_to_text(table_data)
                                text_parts.append(f"\n[Tablo]\n{table_text}\n")
                            break
            
            if current_section["content"]:
                sections.append(current_section)
            
            content = "\n\n".join(text_parts)
            source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
            
            core_props = doc.core_properties
            metadata = {
                "author": core_props.author or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "subject": core_props.subject or "",
            }
            
            return ProcessedDocument(
                content=content,
                source=source,
                file_type="docx",
                title=core_props.title or Path(source).stem if source != "stream" else "",
                metadata=metadata,
                tables=tables,
                sections=sections,
            )
            
        except Exception as e:
            logger.error(f"DOCX işleme hatası: {e}")
            raise
    
    @staticmethod
    def _extract_table(table: Table) -> List[List[str]]:
        """Tabloyu 2D liste olarak çıkar"""
        data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            data.append(row_data)
        return data
    
    @staticmethod
    def _table_to_text(table: List[List[str]]) -> str:
        """Tabloyu text formatına çevir"""
        rows = []
        for row in table:
            cells = [str(cell).strip() if cell else "" for cell in row]
            rows.append(" | ".join(cells))
        return "\n".join(rows)


class ExcelProcessor:
    """Excel doküman işleyici"""
    
    @staticmethod
    def extract_text(file_path: Union[str, Path, BytesIO]) -> ProcessedDocument:
        """Excel'den text çıkar (tüm sheet'ler)"""
        if not EXCEL_AVAILABLE:
            raise ImportError("Excel işleme için openpyxl gerekli")
        
        try:
            if isinstance(file_path, BytesIO):
                wb = openpyxl.load_workbook(file_path, data_only=True)
            else:
                wb = openpyxl.load_workbook(str(file_path), data_only=True)
            
            text_parts = []
            tables = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n## Sheet: {sheet_name}\n")
                
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            row_data.append(str(cell))
                        else:
                            row_data.append("")
                    
                    if any(cell.strip() for cell in row_data):
                        sheet_data.append(row_data)
                
                if sheet_data:
                    tables.append(sheet_data)
                    for row in sheet_data:
                        text_parts.append(" | ".join(row))
            
            content = "\n".join(text_parts)
            source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
            
            return ProcessedDocument(
                content=content,
                source=source,
                file_type="excel",
                title=Path(source).stem if source != "stream" else "",
                metadata={"sheet_count": len(wb.sheetnames), "sheets": wb.sheetnames},
                tables=tables,
            )
            
        except Exception as e:
            logger.error(f"Excel işleme hatası: {e}")
            raise


class PPTXProcessor:
    """PowerPoint doküman işleyici"""
    
    @staticmethod
    def extract_text(file_path: Union[str, Path, BytesIO]) -> ProcessedDocument:
        """PowerPoint'den text çıkar (tüm slide'lar)"""
        if not PPTX_AVAILABLE:
            raise ImportError("PowerPoint işleme için python-pptx gerekli")
        
        try:
            if isinstance(file_path, BytesIO):
                prs = Presentation(file_path)
            else:
                prs = Presentation(str(file_path))
            
            text_parts = []
            sections = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_title = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.level == 0 and not slide_title:
                                    slide_title = para.text.strip()
                        
                        slide_text.append(text)
                
                if slide_text:
                    text_parts.append(f"\n[Slide {i}]\n" + "\n".join(slide_text))
                    sections.append({
                        "title": slide_title or f"Slide {i}",
                        "content": slide_text
                    })
            
            content = "\n\n".join(text_parts)
            source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
            
            return ProcessedDocument(
                content=content,
                source=source,
                file_type="pptx",
                title=Path(source).stem if source != "stream" else "",
                page_count=len(prs.slides),
                metadata={"slide_count": len(prs.slides)},
                sections=sections,
            )
            
        except Exception as e:
            logger.error(f"PowerPoint işleme hatası: {e}")
            raise


class HTMLProcessor:
    """HTML doküman işleyici"""
    
    @staticmethod
    def extract_text(
        file_path: Union[str, Path, BytesIO],
        encoding: Optional[str] = None
    ) -> ProcessedDocument:
        """HTML'den temiz text çıkar"""
        if not BS4_AVAILABLE:
            raise ImportError("HTML işleme için beautifulsoup4 gerekli")
        
        try:
            if isinstance(file_path, BytesIO):
                content = file_path.read()
                if isinstance(content, bytes):
                    if encoding:
                        html = content.decode(encoding)
                    elif CHARDET_AVAILABLE:
                        detected = chardet.detect(content)
                        html = content.decode(detected.get('encoding', 'utf-8'))
                    else:
                        html = content.decode('utf-8', errors='ignore')
                else:
                    html = content
            else:
                path = Path(file_path)
                if encoding:
                    html = path.read_text(encoding=encoding)
                else:
                    try:
                        html = path.read_text(encoding='utf-8')
                    except UnicodeDecodeError:
                        if CHARDET_AVAILABLE:
                            raw = path.read_bytes()
                            detected = chardet.detect(raw)
                            html = raw.decode(detected.get('encoding', 'utf-8'))
                        else:
                            html = path.read_text(encoding='utf-8', errors='ignore')
            
            soup = BeautifulSoup(html, 'html.parser')
            
            for tag in soup(['script', 'style', 'nav', 'footer', 'header', 
                            'aside', 'noscript', 'iframe', 'form']):
                tag.decompose()
            
            title = soup.title.get_text(strip=True) if soup.title else ""
            
            text_parts = []
            tables = []
            
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                level = int(heading.name[1])
                text_parts.append(f"\n{'#' * level} {heading.get_text(strip=True)}\n")
            
            for table in soup.find_all('table'):
                table_data = []
                for row in table.find_all('tr'):
                    cells = []
                    for cell in row.find_all(['td', 'th']):
                        cells.append(cell.get_text(strip=True))
                    if cells:
                        table_data.append(cells)
                
                if table_data:
                    tables.append(table_data)
                    table.decompose()
            
            text = soup.get_text(separator='\n')
            text = re.sub(r'\n{3,}', '\n\n', text)
            text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
            
            source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
            
            return ProcessedDocument(
                content=text,
                source=source,
                file_type="html",
                title=title,
                tables=tables,
            )
            
        except Exception as e:
            logger.error(f"HTML işleme hatası: {e}")
            raise


class TextProcessor:
    """Plain text doküman işleyici"""
    
    @staticmethod
    def extract_text(
        file_path: Union[str, Path, BytesIO],
        encoding: Optional[str] = None
    ) -> ProcessedDocument:
        """Text dosyasından içerik çıkar"""
        try:
            if isinstance(file_path, BytesIO):
                content = file_path.read()
                if isinstance(content, bytes):
                    if encoding:
                        text = content.decode(encoding)
                    elif CHARDET_AVAILABLE:
                        detected = chardet.detect(content)
                        text = content.decode(detected.get('encoding', 'utf-8'))
                    else:
                        text = content.decode('utf-8', errors='ignore')
                else:
                    text = content
            else:
                path = Path(file_path)
                if encoding:
                    text = path.read_text(encoding=encoding)
                else:
                    try:
                        text = path.read_text(encoding='utf-8')
                    except UnicodeDecodeError:
                        if CHARDET_AVAILABLE:
                            raw = path.read_bytes()
                            detected = chardet.detect(raw)
                            text = raw.decode(detected.get('encoding', 'utf-8'))
                        else:
                            text = path.read_text(encoding='utf-8', errors='ignore')
            
            source = str(file_path) if isinstance(file_path, (str, Path)) else "stream"
            
            return ProcessedDocument(
                content=text.strip(),
                source=source,
                file_type="text",
                title=Path(source).stem if source != "stream" else "",
            )
            
        except Exception as e:
            logger.error(f"Text işleme hatası: {e}")
            raise


class DocumentProcessor:
    """
    Ana doküman işleyici sınıfı.
    
    Desteklenen formatlar:
        - PDF (.pdf)
        - Word (.docx, .doc)
        - Excel (.xlsx, .xls)
        - PowerPoint (.pptx, .ppt)
        - HTML (.html, .htm)
        - Text (.txt, .md, .rst)
    """
    
    PROCESSORS = {
        '.pdf': PDFProcessor.extract_text,
        '.docx': DOCXProcessor.extract_text,
        '.doc': DOCXProcessor.extract_text,
        '.xlsx': ExcelProcessor.extract_text,
        '.xls': ExcelProcessor.extract_text,
        '.pptx': PPTXProcessor.extract_text,
        '.ppt': PPTXProcessor.extract_text,
        '.html': HTMLProcessor.extract_text,
        '.htm': HTMLProcessor.extract_text,
        '.txt': TextProcessor.extract_text,
        '.md': TextProcessor.extract_text,
        '.rst': TextProcessor.extract_text,
    }
    
    @classmethod
    def process(
        cls,
        file_path: Union[str, Path, BytesIO],
        file_type: Optional[str] = None,
        **kwargs
    ) -> ProcessedDocument:
        """
        Dosyayı işle ve ProcessedDocument döndür.
        
        Args:
            file_path: Dosya yolu veya BytesIO stream
            file_type: Dosya türü (örn: '.pdf'). None ise otomatik algıla.
            **kwargs: Processor'a özel parametreler
        
        Returns:
            ProcessedDocument
        """
        if file_type is None:
            if isinstance(file_path, (str, Path)):
                file_type = Path(file_path).suffix.lower()
            else:
                raise ValueError("Stream için file_type belirtilmeli")
        
        file_type = file_type.lower()
        if not file_type.startswith('.'):
            file_type = '.' + file_type
        
        processor = cls.PROCESSORS.get(file_type)
        if processor is None:
            raise ValueError(f"Desteklenmeyen dosya türü: {file_type}")
        
        return processor(file_path, **kwargs)
    
    @classmethod
    def is_supported(cls, file_path: Union[str, Path]) -> bool:
        """Dosya türü destekleniyor mu?"""
        ext = Path(file_path).suffix.lower()
        return ext in cls.PROCESSORS
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Desteklenen uzantıları döndür"""
        return list(cls.PROCESSORS.keys())


def process_file(file_path: Union[str, Path], **kwargs) -> ProcessedDocument:
    """Tek dosya işle (convenience function)"""
    return DocumentProcessor.process(file_path, **kwargs)


def process_directory(
    directory: Union[str, Path],
    recursive: bool = True,
    extensions: Optional[List[str]] = None,
) -> List[ProcessedDocument]:
    """
    Dizindeki tüm desteklenen dosyaları işle.
    
    Args:
        directory: Dizin yolu
        recursive: Alt dizinleri de tara
        extensions: Sadece belirli uzantıları işle (None = hepsi)
    
    Returns:
        ProcessedDocument listesi
    """
    directory = Path(directory)
    if not directory.is_dir():
        raise ValueError(f"Dizin bulunamadı: {directory}")
    
    if extensions is None:
        extensions = DocumentProcessor.get_supported_extensions()
    else:
        extensions = [e.lower() if e.startswith('.') else f'.{e.lower()}' for e in extensions]
    
    documents = []
    glob_pattern = "**/*" if recursive else "*"
    
    for file_path in directory.glob(glob_pattern):
        if not file_path.is_file():
            continue
        
        ext = file_path.suffix.lower()
        if ext not in extensions:
            continue
        
        try:
            doc = DocumentProcessor.process(file_path)
            documents.append(doc)
            logger.info(f"✅ İşlendi: {file_path.name} ({doc.word_count} kelime)")
        except Exception as e:
            logger.warning(f"❌ İşlenemedi: {file_path.name} - {e}")
    
    logger.info(f"Toplam {len(documents)} doküman işlendi")
    return documents


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    
    print("Desteklenen formatlar:")
    for ext in DocumentProcessor.get_supported_extensions():
        print(f"  {ext}")
